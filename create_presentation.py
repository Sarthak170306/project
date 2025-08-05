from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_bara_imambara_presentation():
    # Create a new presentation
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "The Majestic Bara Imambara"
    subtitle.text = "A Historic Marvel of Lucknow\nPresented by: [Your Name]\nDate: 2025"
    
    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)  # Purple
    
    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Introduction"
    content_text = """• Bara Imambara is one of the world's most iconic architectural wonders
• Located in Lucknow, Uttar Pradesh, India
• Built between 1784-1791 under Nawab Asaf-ud-Daula
• Originally constructed as a famine relief project
• Features the world's largest unsupported vaulted roof
• Home to the mysterious Bhool Bhulaiya labyrinth
• A perfect blend of Mughal, Persian, and Awadhi architecture"""
    
    content.text = content_text
    
    # Slide 3: Historical Background
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Historical Background"
    content_text = """• Construction began in 1784 during a severe famine
• Built by Nawab Asaf-ud-Daula to provide employment
• Completed in 1791 after 7 years of construction
• Served as a symbol of royal generosity and humanitarian spirit
• Witnessed the Revolt of 1857 against British rule
• Major restorations in 1900 and 1980
• Now a protected monument under Archaeological Survey of India"""
    
    content.text = content_text
    
    # Slide 4: Architectural Marvel
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Architectural Marvel"
    content_text = """• World's largest unsupported vaulted roof
• Built without any wooden beams or iron supports
• Uses lime plaster and bricks for construction
• Features intricate calligraphic carvings
• Grand arched gateways and courtyards
• Central dome with no supporting pillars
• Unique blend of Mughal, Persian, and Awadhi styles"""
    
    content.text = content_text
    
    # Slide 5: Key Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Architectural Features"
    content_text = """• Central Dome: Towering dome without supporting beams
• Bhool Bhulaiya: Complex labyrinth with 1000+ passageways
• Asfi Mosque: Elegant Mughal-style mosque within complex
• Shahi Baoli: Ingenious stepwell and water management system
• Grand Courtyards: Vast open spaces with green lawns
• Arched Gateways: Majestic entrances with intricate details
• Calligraphic Ornamentation: Beautiful Islamic inscriptions"""
    
    content.text = content_text
    
    # Slide 6: Visitor Information
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Visitor Information"
    content_text = """• Opening Hours: 6:00 AM to 5:00 PM (Closed Mondays)
• Location: Chowk, Lucknow, Uttar Pradesh 226003
• Ticket Prices: ₹50 (Adults), ₹25 (Children), ₹500 (Foreigners)
• Photography: ₹10 (Digital), ₹25 (Video)
• Accessibility: Wheelchair ramps and accessible facilities
• Best Time to Visit: October to March (pleasant weather)
• Nearby Metro: Chowk Metro Station (5-minute walk)"""
    
    content.text = content_text
    
    # Slide 7: Cultural Significance
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Cultural Significance"
    content_text = """• Active site of worship and religious ceremonies
• Major tourist attraction drawing millions annually
• Symbol of Lucknow's rich cultural heritage
• Represents Awadhi architecture and craftsmanship
• Important landmark in Indian history
• UNESCO World Heritage Site candidate
• Preserves traditional Islamic architectural techniques
• Showcases humanitarian values through famine relief origins"""
    
    content.text = content_text
    
    # Slide 8: Engineering Feats
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Engineering Achievements"
    content_text = """• Constructed without modern machinery
• Advanced techniques: interlocking stones, lime mortar
• Precision-cut stone blocks and dry stone masonry
• Innovative vaulted roof construction
• Strategic labyrinth design for defense
• Sophisticated water management system
• Earthquake-resistant construction methods
• Enduring structural integrity for over 230 years"""
    
    content.text = content_text
    
    # Slide 9: Tourism Impact
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Tourism and Economic Impact"
    content_text = """• Major contributor to Lucknow's tourism economy
• Generates employment for local guides and vendors
• Promotes cultural exchange and understanding
• Supports local handicraft and souvenir industries
• Educational destination for architecture students
• Photography and film location opportunities
• Cultural events and festivals venue
• Heritage conservation and restoration projects"""
    
    content.text = content_text
    
    # Slide 10: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    content_text = """• Bara Imambara stands as a testament to human ingenuity
• Perfect blend of architectural beauty and engineering excellence
• Rich historical significance and cultural heritage
• Continues to inspire and amaze visitors worldwide
• Symbol of resilience, humanitarian values, and artistic achievement
• Must-visit destination for anyone interested in Indian heritage
• Preserves the legacy of Awadhi culture and craftsmanship
• A living monument that bridges past and present"""
    
    content.text = content_text
    
    # Slide 11: Thank You
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = "Questions & Discussion\n\nContact: [Your Contact Information]\nEmail: [Your Email]\n\nVisit Bara Imambara to experience its majesty firsthand!"
    
    # Style the thank you slide
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)
    
    # Save the presentation
    prs.save('Bara_Imambara_Presentation.pptx')
    print("PowerPoint presentation created successfully: Bara_Imambara_Presentation.pptx")

if __name__ == "__main__":
    create_bara_imambara_presentation() 