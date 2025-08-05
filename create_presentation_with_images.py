from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_bara_imambara_presentation_with_images():
    # Create a new presentation
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "The Majestic Bara Imambara"
    subtitle.text = "A Historic Marvel of Lucknow\nPresented by: [Your Name]\nDate: 2025"
    
    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)  # Purple
    
    # Add background image if available
    try:
        if os.path.exists("2-bada-imambara-lucknow-up-attr-hero.jpg"):
            slide.shapes.add_picture("2-bada-imambara-lucknow-up-attr-hero.jpg", 
                                   Inches(1), Inches(2), height=Inches(4))
    except:
        pass
    
    # Slide 2: Introduction with Image
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Introduction"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)
    title_frame.paragraphs[0].font.bold = True
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    content_frame = content_box.text_frame
    content_text = """• Bara Imambara is one of the world's most iconic architectural wonders
• Located in Lucknow, Uttar Pradesh, India
• Built between 1784-1791 under Nawab Asaf-ud-Daula
• Originally constructed as a famine relief project
• Features the world's largest unsupported vaulted roof
• Home to the mysterious Bhool Bhulaiya labyrinth
• A perfect blend of Mughal, Persian, and Awadhi architecture"""
    
    content_frame.text = content_text
    content_frame.paragraphs[0].font.size = Pt(16)
    
    # Add image
    try:
        if os.path.exists("Discover-the-Majestic-Bara-Imambara-A-Historic-Marvel-of-Lucknow.jpg"):
            slide.shapes.add_picture("Discover-the-Majestic-Bara-Imambara-A-Historic-Marvel-of-Lucknow.jpg", 
                                   Inches(5.5), Inches(2), height=Inches(4))
    except:
        pass
    
    # Slide 3: Historical Background
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Historical Background"
    content_text = """• Construction began in 1784 during a severe famine
• Built by Nawab Asaf-ud-Daula to provide employment
• Completed in 1791 after 7 years of construction
• Served as a symbol of royal generosity and humanitarian spirit
• Witnessed the Revolt of 1857 against British rule
• Major restorations in 1900 and 1980
• Now a protected monument under Archaeological Survey of India"""
    
    content.text = content_text
    
    # Slide 4: Architectural Marvel with Image
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Architectural Marvel"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)
    title_frame.paragraphs[0].font.bold = True
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    content_frame = content_box.text_frame
    content_text = """• World's largest unsupported vaulted roof
• Built without any wooden beams or iron supports
• Uses lime plaster and bricks for construction
• Features intricate calligraphic carvings
• Grand arched gateways and courtyards
• Central dome with no supporting pillars
• Unique blend of Mughal, Persian, and Awadhi styles"""
    
    content_frame.text = content_text
    content_frame.paragraphs[0].font.size = Pt(16)
    
    # Add image
    try:
        if os.path.exists("Bara-Imambara-Lucknow-Tourist-Attraction-1024x683.jpg"):
            slide.shapes.add_picture("Bara-Imambara-Lucknow-Tourist-Attraction-1024x683.jpg", 
                                   Inches(5.5), Inches(2), height=Inches(4))
    except:
        pass
    
    # Slide 5: Key Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Architectural Features"
    content_text = """• Central Dome: Towering dome without supporting beams
• Bhool Bhulaiya: Complex labyrinth with 1000+ passageways
• Asfi Mosque: Elegant Mughal-style mosque within complex
• Shahi Baoli: Ingenious stepwell and water management system
• Grand Courtyards: Vast open spaces with green lawns
• Arched Gateways: Majestic entrances with intricate details
• Calligraphic Ornamentation: Beautiful Islamic inscriptions"""
    
    content.text = content_text
    
    # Slide 6: Visitor Information with Image
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Visitor Information"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)
    title_frame.paragraphs[0].font.bold = True
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    content_frame = content_box.text_frame
    content_text = """• Opening Hours: 6:00 AM to 5:00 PM (Closed Mondays)
• Location: Chowk, Lucknow, Uttar Pradesh 226003
• Ticket Prices: ₹50 (Adults), ₹25 (Children), ₹500 (Foreigners)
• Photography: ₹10 (Digital), ₹25 (Video)
• Accessibility: Wheelchair ramps and accessible facilities
• Best Time to Visit: October to March (pleasant weather)
• Nearby Metro: Chowk Metro Station (5-minute walk)"""
    
    content_frame.text = content_text
    content_frame.paragraphs[0].font.size = Pt(16)
    
    # Add image
    try:
        if os.path.exists("54522.jpg"):
            slide.shapes.add_picture("54522.jpg", 
                                   Inches(5.5), Inches(2), height=Inches(4))
    except:
        pass
    
    # Slide 7: Cultural Significance
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Cultural Significance"
    content_text = """• Active site of worship and religious ceremonies
• Major tourist attraction drawing millions annually
• Symbol of Lucknow's rich cultural heritage
• Represents Awadhi architecture and craftsmanship
• Important landmark in Indian history
• UNESCO World Heritage Site candidate
• Preserves traditional Islamic architectural techniques
• Showcases humanitarian values through famine relief origins"""
    
    content.text = content_text
    
    # Slide 8: Engineering Feats with Image
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Engineering Achievements"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)
    title_frame.paragraphs[0].font.bold = True
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    content_frame = content_box.text_frame
    content_text = """• Constructed without modern machinery
• Advanced techniques: interlocking stones, lime mortar
• Precision-cut stone blocks and dry stone masonry
• Innovative vaulted roof construction
• Strategic labyrinth design for defense
• Sophisticated water management system
• Earthquake-resistant construction methods
• Enduring structural integrity for over 230 years"""
    
    content_frame.text = content_text
    content_frame.paragraphs[0].font.size = Pt(16)
    
    # Add image
    try:
        if os.path.exists("images.jpg"):
            slide.shapes.add_picture("images.jpg", 
                                   Inches(5.5), Inches(2), height=Inches(4))
    except:
        pass
    
    # Slide 9: Tourism Impact
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Tourism and Economic Impact"
    content_text = """• Major contributor to Lucknow's tourism economy
• Generates employment for local guides and vendors
• Promotes cultural exchange and understanding
• Supports local handicraft and souvenir industries
• Educational destination for architecture students
• Photography and film location opportunities
• Cultural events and festivals venue
• Heritage conservation and restoration projects"""
    
    content.text = content_text
    
    # Slide 10: Conclusion with Image
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)
    title_frame.paragraphs[0].font.bold = True
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    content_frame = content_box.text_frame
    content_text = """• Bara Imambara stands as a testament to human ingenuity
• Perfect blend of architectural beauty and engineering excellence
• Rich historical significance and cultural heritage
• Continues to inspire and amaze visitors worldwide
• Symbol of resilience, humanitarian values, and artistic achievement
• Must-visit destination for anyone interested in Indian heritage
• Preserves the legacy of Awadhi culture and craftsmanship
• A living monument that bridges past and present"""
    
    content_frame.text = content_text
    content_frame.paragraphs[0].font.size = Pt(16)
    
    # Add image
    try:
        if os.path.exists("download.jpg"):
            slide.shapes.add_picture("download.jpg", 
                                   Inches(5.5), Inches(2), height=Inches(4))
    except:
        pass
    
    # Slide 11: Thank You
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = "Questions & Discussion\n\nContact: [Your Contact Information]\nEmail: [Your Email]\n\nVisit Bara Imambara to experience its majesty firsthand!"
    
    # Style the thank you slide
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 0, 130)
    
    # Save the presentation
    prs.save('Bara_Imambara_Presentation_With_Images.pptx')
    print("PowerPoint presentation with images created successfully: Bara_Imambara_Presentation_With_Images.pptx")

if __name__ == "__main__":
    create_bara_imambara_presentation_with_images() 